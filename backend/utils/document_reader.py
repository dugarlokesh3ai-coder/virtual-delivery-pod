from __future__ import annotations

import base64
import io
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
from openai import OpenAI

load_dotenv()

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
OCR_MODEL = os.getenv("OCR_MODEL", "gpt-4.1-mini")
MAX_OCR_PDF_PAGES = int(os.getenv("MAX_OCR_PDF_PAGES", "6"))
MAX_IMAGE_BYTES = int(os.getenv("MAX_IMAGE_BYTES", "20000000"))

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
TEXT_EXTENSIONS = {".txt", ".md", ".csv"}


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _mime_type(extension: str) -> str:
    mapping = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".tif": "image/tiff",
        ".tiff": "image/tiff",
    }
    return mapping.get(extension.lower(), "image/png")


def _ocr_image_bytes(image_bytes: bytes, extension: str, context: str = "uploaded image") -> str:
    if not os.getenv("OPENAI_API_KEY"):
        return "[OCR unavailable: OPENAI_API_KEY is not configured.]"

    if len(image_bytes) > MAX_IMAGE_BYTES:
        return f"[OCR skipped: image is larger than {MAX_IMAGE_BYTES} bytes.]"

    encoded = base64.b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{_mime_type(extension)};base64,{encoded}"

    response = client.chat.completions.create(
        model=OCR_MODEL,
        messages=[
            {
                "role": "system",
                "content": (
                    "You extract text and requirement-relevant content from images for a ServiceNow delivery analysis. "
                    "Return concise plain text. Include visible headings, labels, tables, process steps, screenshots, form fields, "
                    "business rules, states, error messages, and any readable notes. Do not invent missing text."
                ),
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": f"Extract all useful text and requirement details from this {context}.",
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": data_url},
                    },
                ],
            },
        ],
        temperature=0,
    )

    return response.choices[0].message.content or ""


def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document
    except Exception as error:
        return f"[DOCX extraction unavailable: python-docx is not installed. {error}]"

    document = Document(io.BytesIO(data))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception as error:
        return f"[PPTX extraction unavailable: python-pptx is not installed. {error}]"

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_parts.append(text.strip())

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        slide_parts.append(" | ".join(cells))

        if slide_parts:
            parts.append(f"Slide {slide_index}:\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def _convert_ppt_to_pptx_if_possible(data: bytes, filename: str) -> Optional[bytes]:
    """Best-effort legacy .ppt support. Requires LibreOffice/soffice on the backend image."""
    with tempfile.TemporaryDirectory() as temp_dir:
        source = Path(temp_dir) / filename
        source.write_bytes(data)

        try:
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    temp_dir,
                    str(source),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=45,
            )
        except Exception:
            return None

        converted = source.with_suffix(".pptx")
        if converted.exists():
            return converted.read_bytes()

    return None


def _extract_pdf_text(data: bytes) -> str:
    text_parts: list[str] = []

    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        for page_index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                text_parts.append(f"Page {page_index}:\n{text.strip()}")
    except Exception as error:
        text_parts.append(f"[PDF text extraction warning: {error}]")

    return "\n\n".join(text_parts)


def _extract_pdf_with_ocr(data: bytes) -> str:
    try:
        import fitz  # PyMuPDF
    except Exception as error:
        return f"[Image-based PDF OCR unavailable: PyMuPDF is not installed. {error}]"

    parts: list[str] = []
    pdf = fitz.open(stream=data, filetype="pdf")
    page_count = min(len(pdf), MAX_OCR_PDF_PAGES)

    for page_index in range(page_count):
        page = pdf[page_index]
        pix = page.get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False)
        image_bytes = pix.tobytes("png")
        extracted = _ocr_image_bytes(
            image_bytes,
            ".png",
            context=f"PDF page {page_index + 1}",
        )
        if extracted.strip():
            parts.append(f"OCR Page {page_index + 1}:\n{extracted.strip()}")

    if len(pdf) > page_count:
        parts.append(f"[OCR limited to first {page_count} pages out of {len(pdf)}.]")

    return "\n\n".join(parts)


def _extract_pdf(data: bytes) -> str:
    text = _extract_pdf_text(data)
    clean_text_length = len("".join(text.split()))

    # Use OCR only when normal text extraction is weak or unavailable.
    if clean_text_length >= 300:
        return text

    ocr_text = _extract_pdf_with_ocr(data)
    if text.strip():
        return f"{text}\n\nImage/OCR Extraction:\n{ocr_text}"

    return ocr_text


async def extract_text_from_file(file) -> str:
    filename = file.filename or "uploaded_file"
    extension = Path(filename).suffix.lower()
    data = await file.read()

    if not data:
        return ""

    try:
        if extension in TEXT_EXTENSIONS:
            return _decode_text(data)

        if extension == ".docx":
            return _extract_docx(data)

        if extension == ".pdf":
            return _extract_pdf(data)

        if extension == ".pptx":
            return _extract_pptx(data)

        if extension == ".ppt":
            converted = _convert_ppt_to_pptx_if_possible(data, filename)
            if converted:
                return _extract_pptx(converted)
            return (
                "[Legacy .ppt extraction requires LibreOffice/soffice on the backend. "
                "Please upload .pptx for best results.]"
            )

        if extension in IMAGE_EXTENSIONS:
            return _ocr_image_bytes(data, extension, context=filename)

        return f"[Unsupported file type: {extension}. Supported: txt, pdf, docx, pptx, ppt, png, jpg, jpeg, webp, gif, bmp, tif, tiff.]"
    except Exception as error:
        return f"[Failed to extract text from {filename}: {error}]"
